import os
import sys
import threading
import webbrowser
from pathlib import Path

from flask import Flask, jsonify, render_template, request, send_from_directory

from app.database import (
    get_connection,
    create_recording,
    delete_recording,
    delete_all_recordings,
    get_all_recordings,
    get_recording,
    init_db,
    update_title_and_memo,
    update_transcript_and_summary,
    update_summary_error,
    update_cleaning_error,
    update_transcript_error,
    update_recording_category,
    get_all_categories,
    create_category,
    update_category,
    delete_category,
    delete_recordings_by_category,
    get_all_vault_memos,
    create_vault_memo,
    update_vault_memo,
    delete_vault_memo,
    get_all_vault_files,
    create_vault_file,
    update_vault_file_summary,
    delete_vault_file,
    delete_vault_items_by_category,
    # Phase 3
    get_all_persona_settings,
    update_persona_enabled,
    create_persona,
    update_persona,
    delete_persona,
    create_council_session,
    update_council_session_decision,
    create_council_adopted,
    get_council_sessions,
    update_adopted_rating,
    get_highly_rated_answers,
    delete_council_session,
    get_context_for_category,
    get_context_for_categories,
    # グループ
    get_all_persona_groups,
    create_persona_group,
    update_persona_group,
    delete_persona_group,
    add_persona_to_default_group,
    remove_persona_from_all_groups,
)
from app.services.recorder import delete_recording as delete_wav
from app.services.recorder import get_status, list_recordings, start, stop

# ── PyInstaller対応：リソースパスの解決 ──────────────
# --onefile で固めた場合、リソースは一時展開ディレクトリに置かれる
# 通常実行時は main.py のある場所をベースにする
if getattr(sys, "frozen", False):
    # PyInstallerで固めた実行ファイルとして起動している
    BASE_DIR = Path(sys.executable).resolve().parent
    RESOURCE_DIR = Path(sys._MEIPASS)  # 一時展開ディレクトリ
else:
    # 通常のPython実行
    BASE_DIR     = Path(__file__).resolve().parent
    RESOURCE_DIR = BASE_DIR

app = Flask(
    __name__,
    template_folder=str(RESOURCE_DIR / "app" / "templates"),
    static_folder=str(RESOURCE_DIR / "app" / "static"),
)

# ── 起動時にDBを初期化 ────────────────────────────
init_db()


# ══════════════════════════════════════════════════
#  ページルーティング
# ══════════════════════════════════════════════════

@app.route("/vault")
def vault():
    """保管庫画面"""
    return render_template("vault.html", active_page="vault")


@app.route("/council")
def council():
    """相談室画面"""
    return render_template("council.html", active_page="council")


@app.route("/inquiry")
def inquiry():
    """照会室画面"""
    return render_template("inquiry.html", active_page="inquiry")


@app.route("/api/inquiry/ask", methods=["POST"])
def inquiry_ask():
    """
    照会室：アーカイバーがSSEでストリーミング回答する
    Request JSON: {
        "question": str,
        "category_id": int,
        "context_keys": [...],
        "history": [{"role": "user"|"assistant", "content": str}, ...]
    }
    """
    import json as _json
    import urllib.request as _ureq
    from app.services.transcriber import _get_ollama_model
    from flask import Response, stream_with_context

    data         = request.get_json(silent=True) or {}
    question     = data.get("question", "").strip()
    # 複数カテゴリ対応：category_ids（リスト）を優先、なければcategory_id（単数）を使用
    cat_ids_raw  = data.get("category_ids")
    if cat_ids_raw:
        category_ids = [int(c) for c in cat_ids_raw if c]
    else:
        category_ids = [int(data.get("category_id", 1) or 1)]
    context_keys = set(data.get("context_keys", None) or [])
    history      = data.get("history", [])  # 会話履歴
    use_all_context = data.get("context_keys") is None

    if not question:
        return jsonify({"status": "error", "message": "照会内容を入力してください"}), 400

    # コンテキスト収集（複数カテゴリ対応）
    ctx = get_context_for_categories(category_ids)

    def build_context_text():
        if not use_all_context and not context_keys:
            return ""

        all_cats = get_all_categories()
        cat_name_map = {c['id']: c['name'] for c in all_cats}
        multi_cat = len(category_ids) > 1

        def cat_label(item):
            if not multi_cat:
                return ""
            cid = item.get('category_id')
            return f"({cat_name_map.get(cid, f'カテゴリ{cid}')})" if cid else ""

        parts = []
        for r in ctx.get("recordings", []):
            key = f"rec_{r['id']}"
            if not use_all_context and key not in context_keys:
                continue
            # ai_summaryを優先、なければcleaned_transcript（5000文字上限）にフォールバック
            ref_text = r.get("ai_summary") or r.get("cleaned_transcript") or ""
            if ref_text:
                parts.append(f"[録音:{r['title']}{cat_label(r)}]\n{ref_text[:5000]}")
        for m in ctx.get("memos", []):
            key = f"memo_{m['id']}"
            if not use_all_context and key not in context_keys:
                continue
            parts.append(f"[メモ:{m['title']}{cat_label(m)}]\n{m['body'][:5000]}")
        for f in ctx.get("files", []):
            key = f"file_{f['id']}"
            if not use_all_context and key not in context_keys:
                continue
            if f.get("summary"):
                parts.append(f"[ファイル:{f['original_name']}{cat_label(f)}]\n{f['summary'][:5000]}")
        return "\n\n".join(parts)

    context_text = build_context_text()

    # アーカイバーのプロンプト取得
    with get_connection() as conn:
        row = conn.execute(
            "SELECT system_prompt FROM persona_settings WHERE persona_name='アーカイバー' LIMIT 1"
        ).fetchone()
    archiver_prompt = row["system_prompt"] if row else ""

    def generate():
        ollama_res = None
        try:
            env     = _read_env()
            ai_mode = env.get("AI_MODE", "ollama")

            from app.services.llm import get_provider, is_local_mode
            provider = get_provider(ai_mode)
            if provider is None:
                yield f"data: {_json.dumps({'type': 'error', 'message': f'不明なAIモード: {ai_mode}'}, ensure_ascii=False)}\n\n"
                return

            api_key = ""
            if not is_local_mode(ai_mode):
                api_key_name = getattr(provider, "API_KEY_NAME", "")
                api_key      = env.get(api_key_name, "").strip() if api_key_name else ""
                if not api_key:
                    yield f"data: {_json.dumps({'type': 'error', 'message': f'{provider.DISPLAY_NAME} のAPIキーが設定されていません'}, ensure_ascii=False)}\n\n"
                    return

            system_content = archiver_prompt
            if context_text:
                system_content += f"\n\n【参照データ】\n{context_text}"

            messages = [{"role": "system", "content": system_content}]
            for h in history:
                role    = h.get("role", "user")
                content = h.get("content", "")
                if role in ("user", "assistant") and content:
                    messages.append({"role": role, "content": content})
            messages.append({"role": "user", "content": question})

            from app.services.transcriber import _get_summary_num_ctx
            options     = {"num_predict": 4096, "num_ctx": _get_summary_num_ctx(), "temperature": 0.2, "repeat_penalty": 1.2}
            full_answer = ""

            if is_local_mode(ai_mode):
                for chunk in provider.chat_stream(messages, options=options):
                    full_answer += chunk
                    yield f"data: {_json.dumps({'type': 'chunk', 'chunk': chunk}, ensure_ascii=False)}\n\n"
            else:
                for chunk in provider.chat_stream(messages, api_key=api_key, options=options):
                    full_answer += chunk
                    yield f"data: {_json.dumps({'type': 'chunk', 'chunk': chunk}, ensure_ascii=False)}\n\n"

            yield f"data: {_json.dumps({'type': 'done', 'answer': full_answer}, ensure_ascii=False)}\n\n"

        except GeneratorExit:
            if ollama_res:
                try: ollama_res.close()
                except Exception: pass
            return
        except Exception as e:
            err_msg = str(e)
            print(f"[inquiry] エラー: {err_msg}")
            yield f"data: {_json.dumps({'type': 'error', 'message': err_msg}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.route("/api/inquiry/feedback", methods=["POST"])
def inquiry_feedback():
    """照会回答へのフィードバック（👍👎）を記録する"""
    data     = request.get_json(silent=True) or {}
    question = data.get("question", "").strip()
    answer   = data.get("answer",   "").strip()
    feedback = data.get("feedback", "")  # 'good' or 'bad'

    if feedback not in ("good", "bad"):
        return jsonify({"status": "error", "message": "feedback は good または bad"}), 400

    print(f"[inquiry] フィードバック: {feedback} | Q: {question[:50]}")
    # 将来DBに保存する場合はここに追加
    return jsonify({"status": "ok"}), 200

@app.route("/settings")
def settings():
    """設定画面"""
    return render_template("settings.html", active_page="settings")


# ══════════════════════════════════════════════════
#  カテゴリ API
# ══════════════════════════════════════════════════

@app.route("/api/categories", methods=["GET"])
def categories_get():
    """カテゴリ一覧を返す"""
    return jsonify(get_all_categories()), 200


@app.route("/api/categories", methods=["POST"])
def categories_create():
    """カテゴリを新規作成する"""
    data  = request.get_json(silent=True) or {}
    name  = data.get("name", "").strip()
    color = data.get("color", "#6B7280").strip()
    if not name:
        return jsonify({"status": "error", "message": "カテゴリ名は必須です"}), 400
    result = create_category(name, color)
    return jsonify(result), 200 if result["status"] == "ok" else 400


@app.route("/api/categories/<int:category_id>", methods=["PUT"])
def categories_update(category_id):
    """カテゴリ名・色を更新する"""
    data  = request.get_json(silent=True) or {}
    name  = data.get("name", "").strip()
    color = data.get("color", "#6B7280").strip()
    if not name:
        return jsonify({"status": "error", "message": "カテゴリ名は必須です"}), 400
    result = update_category(category_id, name, color)
    return jsonify(result), 200 if result["status"] == "ok" else 400


@app.route("/api/categories/<int:category_id>", methods=["DELETE"])
def categories_delete(category_id):
    """カテゴリを削除しデータを未分類に移動する"""
    result = delete_category(category_id)
    return jsonify(result), 200 if result["status"] == "ok" else 400


@app.route("/api/categories/<int:category_id>/recordings", methods=["DELETE"])
def categories_delete_recordings(category_id):
    """指定カテゴリの録音データを全削除する（後方互換）"""
    count = delete_recordings_by_category(category_id)
    return jsonify({"status": "ok", "deleted": count}), 200


@app.route("/api/categories/<int:category_id>/data", methods=["DELETE"])
def categories_delete_all_data(category_id):
    """指定カテゴリの録音・メモ・ファイル・Webデータをすべて削除する"""
    import os
    deleted = {"recordings": 0, "memos": 0, "files": 0}

    # 録音データ削除
    all_records = get_all_recordings(category_id=category_id)
    for r in all_records:
        if r.get("wav_file"):
            delete_wav(r["wav_file"])
        delete_recording(r["id"])
        deleted["recordings"] += 1

    # メモ削除
    with get_connection() as conn:
        rows = conn.execute("SELECT id FROM vault_memos WHERE category_id=?", (category_id,)).fetchall()
        for row in rows:
            conn.execute("DELETE FROM vault_memos WHERE id=?", (row["id"],))
            deleted["memos"] += 1
        # ファイル削除
        frows = conn.execute("SELECT id, filename FROM vault_files WHERE category_id=?", (category_id,)).fetchall()
        for row in frows:
            fpath = RESOURCE_DIR / "uploads" / row["filename"]
            if fpath.exists():
                fpath.unlink()
            conn.execute("DELETE FROM vault_files WHERE id=?", (row["id"],))
            deleted["files"] += 1
        conn.commit()

    total = sum(deleted.values())
    print(f"[data] カテゴリ{category_id}のデータ削除: {deleted}")
    return jsonify({"status": "ok", "deleted": total, "detail": deleted}), 200


@app.route("/api/recordings/<int:record_id>/category", methods=["PUT"])
def recording_update_category(record_id):
    """録音のカテゴリを更新する"""
    data        = request.get_json(silent=True) or {}
    category_id = data.get("category_id", 1)
    result      = update_recording_category(record_id, category_id)
    return jsonify(result), 200

@app.route("/api/shutdown", methods=["POST"])
def shutdown():
    """COVEを終了する"""
    import threading
    def _shutdown():
        import time, os, signal, subprocess, platform
        time.sleep(0.3)
        print("[main] シャットダウン要求を受信しました")
        try:
            if platform.system() == "Windows":
                subprocess.run(["taskkill", "/F", "/IM", "ollama.exe"], capture_output=True)
            else:
                # Mac/Linux: pkill or killall
                subprocess.run(["pkill", "-f", "ollama"], capture_output=True)
            print("[main] Ollamaを停止しました")
        except Exception as e:
            print(f"[main] Ollama停止エラー: {e}")
        os.kill(os.getpid(), signal.SIGTERM)
    threading.Thread(target=_shutdown, daemon=True).start()
    return jsonify({"status": "ok"}), 200

@app.route("/api/disk/status", methods=["GET"])
def disk_status():
    """ディスクの空き容量を返す"""
    import shutil
    try:
        usage = shutil.disk_usage(str(BASE_DIR))
        return jsonify({
            "total": usage.total,
            "used":  usage.used,
            "free":  usage.free,
            "free_gb":  round(usage.free  / (1024**3), 1),
            "total_gb": round(usage.total / (1024**3), 1),
            "percent_used": round(usage.used / usage.total * 100, 1),
        }), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/system/ram", methods=["GET"])
def system_ram():
    """利用可能なRAM情報を返す"""
    try:
        import psutil
        mem = psutil.virtual_memory()
        total_gb  = mem.total / (1024 ** 3)
        usable_gb = round(total_gb * 0.8, 1)

        ctx_table = [
            {"minutes": 30,  "num_ctx": 16384,  "required_gb": 8.0},
            {"minutes": 60,  "num_ctx": 32768,  "required_gb": 16.0},
            {"minutes": 90,  "num_ctx": 65536,  "required_gb": 32.0},
            {"minutes": 120, "num_ctx": 131072, "required_gb": 56.0},
        ]
        options = [e for e in ctx_table if e["required_gb"] <= usable_gb]
        if not options:
            options = [ctx_table[0]]

        return jsonify({
            "total_gb":  round(total_gb, 1),
            "avail_gb":  round(mem.available / (1024 ** 3), 1),
            "usable_gb": usable_gb,
            "options":   options,
        }), 200
    except ImportError:
        return jsonify({
            "total_gb": 0, "avail_gb": 0, "usable_gb": 0,
            "options": [
                {"minutes": 30,  "num_ctx": 16384,  "required_gb": 8.0},
                {"minutes": 60,  "num_ctx": 32768,  "required_gb": 16.0},
            ],
        }), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/help")
def help_page():
    return render_template("help.html")


@app.route("/favicon.ico")
def favicon():
    return send_from_directory(
        str(BASE_DIR / "app" / "static"),
        "cove.ico",
        mimetype="image/x-icon"
    )


@app.route("/")
def index():
    """メイン画面（録音・一覧）"""
    return render_template("index.html", active_page="index")




# ══════════════════════════════════════════════════
#  保管庫 API
# ══════════════════════════════════════════════════

VAULT_DIR = BASE_DIR / "vault"
VAULT_DIR.mkdir(exist_ok=True)
_file_summarize_status   = {"status": "idle", "message": ""}
_recording_upload_status = {"status": "idle", "message": "", "progress": 0, "step": ""}
_transcribe_progress     = {"status": "idle", "step": "", "message": "", "progress": 0, "elapsed": 0, "record_id": None}
_transcribe_start_time   = 0

@app.route("/api/transcribe/progress", methods=["GET"])
def transcribe_progress():
    """文字起こし・クリーニング・要約の進捗を返す"""
    p = dict(_transcribe_progress)
    if _transcribe_start_time > 0:
        import time as _time
        p["elapsed"] = int(_time.time() - _transcribe_start_time)
    return jsonify(p), 200

# ── Ollama処理の中断管理 ──────────────────────────────
import threading as _threading
_ollama_abort_events: dict[str, "_threading.Event"] = {}  # task_key -> Event

def _new_abort_event(task_key: str) -> "_threading.Event":
    """新しい中断イベントを作成・登録して返す"""
    ev = _threading.Event()
    _ollama_abort_events[task_key] = ev
    return ev

def abort_ollama_task(task_key: str):
    """指定タスクの中断イベントをセットする"""
    ev = _ollama_abort_events.get(task_key)
    if ev:
        ev.set()
        print(f"[ollama] 中断リクエスト: {task_key}")

@app.route("/api/ollama/abort", methods=["POST"])
def ollama_abort():
    """実行中のOllama処理を中断する"""
    data     = request.get_json(silent=True) or {}
    task_key = data.get("task_key", "")
    if task_key:
        abort_ollama_task(task_key)
    else:
        # task_key未指定なら全タスクを中断
        for key in list(_ollama_abort_events.keys()):
            abort_ollama_task(key)
    return jsonify({"status": "ok"}), 200


@app.route("/api/vault/files/status", methods=["GET"])
def vault_file_status():
    return jsonify(_file_summarize_status), 200

@app.route("/api/vault/recording/status", methods=["GET"])
def vault_recording_status():
    return jsonify(_recording_upload_status), 200

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"}
ALLOWED_AUDIO_EXTENSIONS = {".wav", ".mp3", ".m4a", ".ogg", ".flac", ".webm"}


@app.route("/api/vault/recording/upload", methods=["POST"])
def vault_recording_upload():
    """保管庫から録音ファイルをアップロードして文字起こし・要約を実施する"""
    global _recording_upload_status

    if "file" not in request.files:
        return jsonify({"status": "error", "message": "ファイルが選択されていません"}), 400

    file        = request.files["file"]
    title       = request.form.get("title", "").strip() or file.filename
    memo        = request.form.get("memo", "").strip()
    category_id = int(request.form.get("category_id", 1) or 1)

    suffix = Path(file.filename).suffix.lower()
    if suffix not in ALLOWED_AUDIO_EXTENSIONS:
        return jsonify({"status": "error", "message": f"対応形式: {', '.join(ALLOWED_AUDIO_EXTENSIONS)}"}), 400

    # WAVとして保存（mp3等はffmpegで変換が必要だがまずwavで保存）
    import uuid, shutil
    UPLOADS_DIR.mkdir(exist_ok=True)
    unique_name = f"vault_{uuid.uuid4().hex}{suffix}"
    save_path   = UPLOADS_DIR / unique_name
    file.save(str(save_path))

    # mp3/m4a等はwavに変換
    wav_path = save_path
    if suffix != ".wav":
        try:
            import subprocess
            wav_name = save_path.with_suffix(".wav").name
            wav_path = UPLOADS_DIR / wav_name
            subprocess.run(
                ["ffmpeg", "-y", "-i", str(save_path), str(wav_path)],
                check=True, capture_output=True, timeout=None
            )
            save_path.unlink(missing_ok=True)
        except Exception:
            wav_path = save_path  # 変換失敗時はそのまま使用

    # DBにレコード作成
    record_id = create_recording(str(wav_path.name), title, memo, category_id)
    abort_ev  = _new_abort_event(f"rec_{record_id}")

    def _process():
        global _recording_upload_status
        try:
            if abort_ev.is_set():
                _recording_upload_status = {"status": "idle", "message": "中断されました", "progress": 0, "step": ""}
                return
            _recording_upload_status = {"status": "running", "message": "文字起こし中…", "progress": 10, "step": "transcribe"}
            from app.services.transcriber import transcribe_and_summarize
            result = transcribe_and_summarize(wav_path.name, "", record_id, abort_event=abort_ev)

            if abort_ev.is_set():
                _recording_upload_status = {"status": "idle", "message": "中断されました", "progress": 0, "step": ""}
                return

            if result["status"] == "error":
                from app.services.transcriber import _parse_api_error as _pae
                err_msg = result.get('message', '不明なエラー')
                update_transcript_error(record_id, err_msg)
                _recording_upload_status = {
                    "status": "error",
                    "message": err_msg,
                    "progress": 0, "step": "error"
                }
                return

            _recording_upload_status = {"status": "running", "message": "クリーニング・要約完了", "progress": 90, "step": "summary"}

            update_transcript_and_summary(record_id, result["transcript"], result["ai_summary"])
            if result.get("cleaned_transcript"):
                from app.database import update_cleaned_transcript
                update_cleaned_transcript(record_id, result["cleaned_transcript"])

            _recording_upload_status = {
                "status": "done",
                "message": f"完了：{title}",
                "progress": 100,
                "step": "done",
                "record_id": record_id,
            }

        except Exception as e:
            _recording_upload_status = {
                "status": "error",
                "message": f"処理中にエラーが発生しました: {str(e)}",
                "progress": 0, "step": "error"
            }

    _recording_upload_status = {"status": "running", "message": "アップロード完了。処理を開始します…", "progress": 5, "step": "start"}
    threading.Thread(target=_process, daemon=True).start()
    return jsonify({"status": "ok", "record_id": record_id}), 200


@app.route("/api/vault/memos", methods=["GET"])
def vault_memos_get():
    category_id = request.args.get("category_id", type=int)
    return jsonify(get_all_vault_memos(category_id=category_id)), 200


@app.route("/api/vault/memos", methods=["POST"])
def vault_memos_create():
    data        = request.get_json(silent=True) or {}
    title       = data.get("title", "").strip()
    body        = data.get("body", "").strip()
    category_id = int(data.get("category_id", 1) or 1)
    if not body:
        return jsonify({"status": "error", "message": "本文は必須です"}), 400
    memo_id = create_vault_memo(title or "無題", body, category_id)
    return jsonify({"status": "ok", "id": memo_id}), 200


@app.route("/api/vault/memos/<int:memo_id>", methods=["PUT"])
def vault_memos_update(memo_id):
    data        = request.get_json(silent=True) or {}
    title       = data.get("title", "").strip()
    body        = data.get("body", "").strip()
    category_id = int(data.get("category_id", 1) or 1)
    if not body:
        return jsonify({"status": "error", "message": "本文は必須です"}), 400
    return jsonify(update_vault_memo(memo_id, title or "無題", body, category_id)), 200


@app.route("/api/vault/memos/<int:memo_id>", methods=["DELETE"])
def vault_memos_delete(memo_id):
    return jsonify(delete_vault_memo(memo_id)), 200


@app.route("/api/vault/files", methods=["GET"])
def vault_files_get():
    category_id = request.args.get("category_id", type=int)
    return jsonify(get_all_vault_files(category_id=category_id)), 200


@app.route("/api/vault/files", methods=["POST"])
def vault_files_upload():
    """ファイルをアップロードしてOllamaで要約する"""
    if "file" not in request.files:
        return jsonify({"status": "error", "message": "ファイルが選択されていません"}), 400

    file        = request.files["file"]
    category_id = int(request.form.get("category_id", 1) or 1)
    ext         = Path(file.filename).suffix.lower()

    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"status": "error", "message": f"対応していないファイル形式です: {ext}"}), 400

    # ファイルを保存
    import uuid
    filename    = f"{uuid.uuid4().hex}{ext}"
    save_path   = VAULT_DIR / filename
    file.save(str(save_path))

    # DBに登録
    file_id = create_vault_file(filename, file.filename, ext, category_id)

    # バックグラウンドでOllama要約
    abort_ev = _new_abort_event(f"file_{file_id}")

    def _summarize():
        global _file_summarize_status
        try:
            _file_summarize_status = {"status": "extracting", "message": f"テキストを抽出中: {file.filename}"}
            text = _extract_text(save_path, ext)
            if not text:
                update_vault_file_summary(file_id, "（テキストを抽出できませんでした）")
                _file_summarize_status = {"status": "done", "message": "完了"}
                return
            if len(text) < 200:
                update_vault_file_summary(file_id, text)
                _file_summarize_status = {"status": "done", "message": f"完了: {file.filename}"}
                return

            if abort_ev.is_set():
                _file_summarize_status = {"status": "idle", "message": "中断されました"}
                return

            _file_summarize_status = {"status": "summarizing", "message": f"LLMで整形中: {file.filename}"}
            from app.services.transcriber import _read_env as _tr_read_env
            _env     = _tr_read_env()
            _ai_mode = _env.get("AI_MODE", "ollama")

            _file_prompt = (
                "以下はファイルから抽出したテキストです。\n"
                "元の内容を一切変えず・追加せず、そのまま読みやすく整理してください。\n"
                "・重要な数値・固有名詞・日付は正確に残す\n"
                "・箇条書きや見出しを使って構造化する\n"
                "・不要な繰り返しや装飾文字のみ除去する\n"
                "・内容の言い換え・要約・追記は絶対にしないこと\n"
                "整理後のテキストのみ出力してください。\n\n"
                f"{text}"
            )

            from app.services.llm import get_provider, is_local_mode
            _provider = get_provider(_ai_mode)
            _messages = [{"role": "user", "content": _file_prompt}]

            if is_local_mode(_ai_mode):
                summary = _provider.generate(_file_prompt)
            else:
                _api_key_name = getattr(_provider, "API_KEY_NAME", "")
                _api_key      = _env.get(_api_key_name, "").strip() if _api_key_name else ""
                if _api_key:
                    summary = _provider.chat(_messages, api_key=_api_key, options={"temperature": 0.0, "max_tokens": 2048})
                else:
                    summary = text  # APIキー未設定時は元テキストをそのまま使用

            update_vault_file_summary(file_id, summary)
            print(f"[vault] ファイル要約完了: {file.filename}")
            _file_summarize_status = {"status": "done", "message": f"完了: {file.filename}"}
        except Exception as e:
            print(f"[vault] ファイル要約エラー: {e}")
            from app.services.transcriber import _parse_api_error as _pae2
            from app.services.llm import get_provider as _gp2
            _env2    = _read_env()
            _mode2   = _env2.get("AI_MODE", "ollama")
            _prov2   = _gp2(_mode2)
            _pname2  = getattr(_prov2, "DISPLAY_NAME", _mode2) if _prov2 else _mode2
            err_msg2 = _pae2(e, _pname2)
            update_vault_file_summary(file_id, f"（{err_msg2}）")
            _file_summarize_status = {"status": "error", "message": err_msg2}

    import threading
    threading.Thread(target=_summarize, daemon=True).start()

    return jsonify({"status": "ok", "id": file_id, "message": "アップロードしました。要約処理中です。"}), 200


def _extract_text(path: Path, ext: str) -> str:
    """ファイルからテキストを抽出する"""
    try:
        if ext == ".pdf":
            import fitz
            doc  = fitz.open(str(path))
            return "\n".join(page.get_text() for page in doc)
        elif ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".xlsx":
            import openpyxl
            wb   = openpyxl.load_workbook(str(path), data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    r = [str(c) for c in row if c is not None]
                    if r:
                        rows.append("\t".join(r))
            return "\n".join(rows)
        elif ext == ".pptx":
            from pptx import Presentation
            prs  = Presentation(str(path))
            rows = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        rows.append(shape.text)
            return "\n".join(rows)
        elif ext in (".txt", ".md"):
            return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        from app.services.transcriber import _parse_api_error as _pae4
        err_msg4 = _pae4(e)
        print(f"[vault] テキスト抽出エラー: {err_msg4}")
    return ""



@app.route("/api/vault/files/<int:file_id>/category", methods=["PUT"])
def vault_file_update_category(file_id):
    """ファイルのカテゴリ・整形内容を更新する"""
    from app.database import get_connection
    data        = request.get_json(silent=True) or {}
    category_id = int(data.get("category_id", 1) or 1)
    summary     = data.get("summary", None)
    with get_connection() as conn:
        if summary is not None:
            conn.execute(
                "UPDATE vault_files SET category_id = ?, summary = ? WHERE id = ?",
                (category_id, summary, file_id)
            )
        else:
            conn.execute(
                "UPDATE vault_files SET category_id = ? WHERE id = ?",
                (category_id, file_id)
            )
        conn.commit()
    return jsonify({"status": "ok"}), 200

@app.route("/api/vault/files/<int:file_id>", methods=["DELETE"])
def vault_files_delete(file_id):
    result = delete_vault_file(file_id)
    if result["status"] == "ok":
        # 実ファイルも削除
        file_path = VAULT_DIR / result["filename"]
        if file_path.exists():
            file_path.unlink()
    return jsonify(result), 200



# Web取得処理の状態管理
_web_fetch_status = {"status": "idle", "message": "", "progress": 0, "total": 0}


@app.route("/api/vault/web/status", methods=["GET"])
def vault_web_status():
    """Web取得処理の状態を返す"""
    return jsonify(_web_fetch_status), 200


@app.route("/api/vault/web/analyze", methods=["POST"])
def vault_web_analyze():
    """URLを解析して同一ドメインのリンク一覧を返す"""
    import urllib.request as _req
    import urllib.parse as _parse
    import html as _html
    import re

    data = request.get_json(silent=True) or {}
    url  = data.get("url", "").strip()
    if not url:
        return jsonify({"status": "error", "message": "URLは必須です"}), 400

    try:
        parsed_base = _parse.urlparse(url)
        base_domain = parsed_base.scheme + "://" + parsed_base.netloc

        headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "ja,en-US;q=0.9,en;q=0.8",
            "Accept-Encoding": "gzip, deflate",
        }
        req      = _req.Request(url, headers=headers)
        with _req.urlopen(req, timeout=15) as res:
            raw      = res.read()
            charset  = res.headers.get_content_charset() or "utf-8"
            html_text = raw.decode(charset, errors="ignore")

        # ページタイトル
        title_match = re.search(r"<title[^>]*>([^<]+)</title>", html_text, re.IGNORECASE)
        page_title  = _html.unescape(title_match.group(1).strip()) if title_match else url

        # 同一ドメインのリンクを抽出
        links = []
        seen  = set()
        for href, text in re.findall(r'href=["\']([^"\']+)["\'][^>]*>([^<]*)', html_text):
            href = href.strip()
            text = _html.unescape(text.strip())

            # 相対URLを絶対URLに変換
            if href.startswith('/'):
                href = base_domain + href
            elif not href.startswith('http'):
                continue

            # 同一ドメインのみ
            parsed = _parse.urlparse(href)
            if parsed.netloc != parsed_base.netloc:
                continue

            # フラグメント・クエリを除いたURLで重複チェック
            clean = parsed.scheme + "://" + parsed.netloc + parsed.path
            if clean in seen or clean == url:
                continue
            seen.add(clean)

            links.append({
                "url":   href,
                "title": text or href,
            })

        return jsonify({
            "status":     "ok",
            "page_title": page_title,
            "base_url":   url,
            "links":      links[:100],  # 最大100件表示
        }), 200

    except Exception as e:
        return jsonify({"status": "error", "message": f"解析エラー: {str(e)}"}), 500


@app.route("/api/vault/web", methods=["POST"])
def vault_web_fetch():
    """選択したURLリストを順番に取得・整形・保管庫に保存する"""
    global _web_fetch_status

    data        = request.get_json(silent=True) or {}
    urls        = data.get("urls", [])
    category_id = int(data.get("category_id", 1) or 1)

    if not urls:
        return jsonify({"status": "error", "message": "URLが選択されていません"}), 400

    def _fetch_all():
        global _web_fetch_status
        import urllib.request as _req
        import urllib.parse as _parse
        import html as _html
        import re

        total = len(urls)
        _web_fetch_status = {"status": "running", "message": f"0 / {total} ページ完了", "progress": 0, "total": total}

        for i, item in enumerate(urls):
            url   = item.get("url", "")
            title = item.get("title", url)
            try:
                _web_fetch_status = {
                    "status":   "running",
                    "message":  f"{i + 1} / {total} ページ取得中: {title[:30]}",
                    "progress": i,
                    "total":    total,
                }

                # ページ取得
                headers = {
                    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
                    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
                    "Accept-Language": "ja,en-US;q=0.9,en;q=0.8",
                    "Accept-Encoding": "gzip, deflate",
                }
                req     = _req.Request(url, headers=headers)
                with _req.urlopen(req, timeout=15) as res:
                    raw     = res.read()
                    charset = res.headers.get_content_charset() or "utf-8"
                    html_text = raw.decode(charset, errors="ignore")

                # タイトル抽出
                title_match = re.search(r"<title[^>]*>([^<]+)</title>", html_text, re.IGNORECASE)
                page_title  = _html.unescape(title_match.group(1).strip()) if title_match else title

                # テキスト抽出
                text = re.sub(r"<script[^>]*>.*?</script>", "", html_text, flags=re.DOTALL|re.IGNORECASE)
                text = re.sub(r"<style[^>]*>.*?</style>",  "", text,      flags=re.DOTALL|re.IGNORECASE)
                text = re.sub(r"<[^>]+>", " ", text)
                text = _html.unescape(text)
                text = re.sub(r"\s+", " ", text).strip()

                if len(text) > 8000:
                    text = text[:8000] + "..."

                if not text:
                    create_vault_memo(page_title, "（テキストを抽出できませんでした）\nURL: " + url, category_id)
                    continue

                # LLMで整形
                _web_fetch_status["message"] = f"{i + 1} / {total} 整形中: {page_title[:30]}"
                from app.services.transcriber import _read_env as _tr_read_env
                _env     = _tr_read_env()
                _ai_mode = _env.get("AI_MODE", "ollama")

                prompt = (
                    "以下はWebページから取得したテキストです。\n"
                    "ナビゲーション・フッター・広告・メニュー等の不要な部分を除き、"
                    "本文の情報を欠如なく整理してください。\n"
                    "箇条書きや見出しを使って読みやすく整形してください。\n"
                    "整形後のテキストのみ出力してください。\n\n"
                    f"{text}"
                )

                from app.services.llm import get_provider, is_local_mode
                _provider = get_provider(_ai_mode)
                _messages = [{"role": "user", "content": prompt}]

                if is_local_mode(_ai_mode):
                    arranged = _provider.generate(prompt)
                else:
                    _api_key_name = getattr(_provider, "API_KEY_NAME", "")
                    _api_key      = _env.get(_api_key_name, "").strip() if _api_key_name else ""
                    if _api_key:
                        arranged = _provider.chat(_messages, api_key=_api_key, options={"temperature": 0.0, "max_tokens": 2048})
                    else:
                        arranged = text

                body = "URL: " + url + "\n\n" + (arranged or text)
                create_vault_memo(page_title, body, category_id)
                print(f"[vault] Web取得完了 ({i+1}/{total}): {page_title}")

            except Exception as e:
                print(f"[vault] Web取得エラー ({url}): {e}")
                from app.services.transcriber import _parse_api_error as _pae3
                err_msg3 = _pae3(e)
                create_vault_memo(
                    "取得失敗: " + (title or url),
                    "URL: " + url + "\n\nエラー: " + err_msg3,
                    category_id
                )

        _web_fetch_status = {
            "status":   "done",
            "message":  f"{total} ページの取得・保存が完了しました",
            "progress": total,
            "total":    total,
        }

    import threading
    threading.Thread(target=_fetch_all, daemon=True).start()
    _web_fetch_status = {"status": "running", "message": "処理を開始しました", "progress": 0, "total": len(urls)}
    return jsonify({"status": "ok", "message": f"{len(urls)}ページの取得を開始しました"}), 200


@app.route("/api/vault/files/<int:file_id>/open-folder", methods=["POST"])
def vault_open_folder(file_id):
    """保存フォルダをエクスプローラーで開く"""
    import subprocess, platform
    if platform.system() == "Windows":
        subprocess.Popen(f'explorer "{VAULT_DIR}"')
    else:
        subprocess.Popen(["open", str(VAULT_DIR)])
    return jsonify({"status": "ok"}), 200




# ══════════════════════════════════════════════════
#  録音 API
# ══════════════════════════════════════════════════

@app.route("/api/record/start", methods=["POST"])
def record_start():
    """
    録音を開始する。
    Response: {"status": "started"} or {"status": "error", "message": str}
    """
    result = start()
    status_code = 200 if result["status"] == "started" else 500
    return jsonify(result), status_code


@app.route("/api/record/stop", methods=["POST"])
def record_stop():
    """
    録音を停止してDBに登録する。
    Request JSON（省略可）: {"title": str, "memo": str}
    Response: {"status": "stopped", "record_id": int, "filename": str, "duration": float}
    """
    result = stop()
    if result["status"] == "error":
        return jsonify(result), 500

    # リクエストからタイトル・メモ・カテゴリを取得
    data        = request.get_json(silent=True) or {}
    title       = data.get("title", "無題").strip() or "無題"
    memo        = data.get("memo", "").strip()
    category_id = int(data.get("category_id", 1) or 1)

    # DBに登録
    record_id = create_recording(
        wav_file=result["filename"],
        title=title,
        memo=memo,
        category_id=category_id,
    )

    return jsonify({
        "status": "stopped",
        "record_id": record_id,
        "filename": result["filename"],
        "duration": result["duration"],
    }), 200


@app.route("/api/record/status", methods=["GET"])
def record_status():
    """
    現在の録音状態を返す。
    Response: {"recording": bool}
    """
    return jsonify(get_status()), 200


# ══════════════════════════════════════════════════
#  文字起こし & 要約 API
# ══════════════════════════════════════════════════

@app.route("/api/transcribe", methods=["POST"])
def transcribe():
    """文字起こしと要約を実行してDBに保存する"""
    global _transcribe_progress, _transcribe_start_time
    import time as _time

    data      = request.get_json(silent=True) or {}
    record_id = data.get("record_id")

    if not record_id:
        return jsonify({"status": "error", "message": "record_id が必要です"}), 400

    record = get_recording(record_id)
    if not record:
        return jsonify({"status": "error", "message": "データが見つかりません"}), 404

    _transcribe_start_time = _time.time()
    _transcribe_progress   = {"status": "running", "step": "transcribe", "message": "🎙️ 文字起こし開始…", "progress": 5, "record_id": record_id}

    def on_progress(step, message, progress, transcript=None, cleaned_transcript=None):
        global _transcribe_progress
        _transcribe_progress = {
            "status":             "running",
            "step":               step,
            "message":            message,
            "progress":           progress,
            "record_id":          record_id,
            "transcript":         transcript or _transcribe_progress.get("transcript", ""),
            "cleaned_transcript": cleaned_transcript or _transcribe_progress.get("cleaned_transcript", ""),
        }

    extra_prompt = data.get("extra_prompt", "").strip()
    from app.services.transcriber import transcribe_and_summarize
    result = transcribe_and_summarize(record["wav_file"], extra_prompt, record_id, progress_callback=on_progress)

    if result["status"] == "error":
        err_msg = result.get("message", "不明なエラーが発生しました")
        if record_id:
            update_transcript_error(record_id, err_msg)
        _transcribe_progress = {"status": "error", "step": "error", "message": f"⚠️ {err_msg}", "progress": 0, "record_id": record_id}
        return jsonify(result), 500

    if result.get("cleaned_transcript"):
        from app.database import update_cleaned_transcript
        update_cleaned_transcript(record_id, result["cleaned_transcript"])

    update_transcript_and_summary(record_id, result["transcript"], result["ai_summary"])
    _transcribe_progress = {"status": "done", "step": "done", "message": "✅ 完了！", "progress": 100, "record_id": record_id}

    return jsonify({
        "status":             "done",
        "transcript":         result["transcript"],
        "cleaned_transcript": result.get("cleaned_transcript", ""),
        "ai_summary":         result["ai_summary"],
        "summary_error":      "",
        "cleaning_error":     "",
    }), 200


# ══════════════════════════════════════════════════
#  録音データ管理 API
# ══════════════════════════════════════════════════

@app.route("/api/recordings", methods=["GET"])
def recordings_list():
    """
    過去の録音データを全件返す。
    Response: [{"id": int, "title": str, ...}, ...]
    """
    category_id = request.args.get("category_id", type=int)
    return jsonify(get_all_recordings(category_id=category_id)), 200


@app.route("/api/recordings/<int:record_id>", methods=["PATCH"])
def recordings_update(record_id: int):
    """
    タイトルと概要メモを更新する。
    Request JSON: {"title": str, "memo": str}
    Response: {"status": "updated"} or {"status": "error", "message": str}
    """
    data = request.get_json(silent=True) or {}
    title = data.get("title", "").strip()
    memo = data.get("memo", "").strip()

    if not title:
        return jsonify({"status": "error", "message": "タイトルは必須です"}), 400

    result = update_title_and_memo(record_id, title, memo)
    status_code = 200 if result["status"] == "updated" else 500
    return jsonify(result), status_code


@app.route("/api/recordings/<int:record_id>", methods=["DELETE"])
def recordings_delete(record_id: int):
    """
    録音データをDBとWAVファイルの両方から削除する。
    Response: {"status": "deleted"} or {"status": "error", "message": str}
    """
    # DBから削除（WAVファイル名も返ってくる）
    result = delete_recording(record_id)
    if result["status"] == "error":
        return jsonify(result), 404

    # WAVファイルも削除
    if result.get("wav_file"):
        delete_wav(result["wav_file"])

    return jsonify({"status": "deleted"}), 200



@app.route("/api/recordings/all", methods=["DELETE"])
def recordings_delete_all():
    """
    すべてのユーザーデータを削除する（録音・メモ・ファイル・Web・相談履歴）
    カテゴリ・ペルソナ設定は維持する
    """
    count = 0

    # ① 録音データ
    all_records = get_all_recordings()
    for record in all_records:
        if record.get("wav_file"):
            delete_wav(record["wav_file"])
        delete_recording(record["id"])
        count += 1

    # ② メモ・ファイル・Web・相談履歴
    with get_connection() as conn:
        # 保管庫ファイルの実体も削除
        frows = conn.execute("SELECT filename FROM vault_files").fetchall()
        for row in frows:
            fpath = RESOURCE_DIR / "uploads" / row["filename"]
            if fpath.exists():
                fpath.unlink()

        rows = conn.execute("DELETE FROM vault_memos").rowcount
        count += rows
        rows = conn.execute("DELETE FROM vault_files").rowcount
        count += rows
        rows = conn.execute("DELETE FROM council_sessions").rowcount
        count += rows
        conn.commit()

    return jsonify({"status": "deleted", "count": count}), 200


# ══════════════════════════════════════════════════
#  音声ファイル配信 API
# ══════════════════════════════════════════════════

if getattr(sys, "frozen", False):
    UPLOADS_DIR = Path(sys.executable).resolve().parent / "uploads"
else:
    UPLOADS_DIR = Path(__file__).resolve().parent / "uploads"

@app.route("/api/audio/<filename>")
def serve_audio(filename):
    """
    WAVファイルをブラウザに配信する（過去データの再生用）。
    ディレクトリトラバーサル対策のため send_from_directory を使用。
    """
    return send_from_directory(str(UPLOADS_DIR), filename)


# ══════════════════════════════════════════════════
#  録音デバイス API
# ══════════════════════════════════════════════════

@app.route("/api/devices", methods=["GET"])
def get_devices():
    """
    利用可能な録音デバイス一覧を返す。
    システム音声デバイス（ステレオミキサー・BlackHole）を自動でハイライト。
    Response: [{"id": int, "name": str, "is_system_audio": bool}, ...]
    """
    import sounddevice as sd
    devices = sd.query_devices()
    result  = []
    # システム音声デバイスのキーワード
    system_keywords = [
        "stereo mix", "ステレオ ミキサー", "ステレオミキサー",
        "blackhole", "black hole",
        "loopback", "what u hear", "wave out mix",
        "virtual audio", "soundflower",
    ]
    for i, dev in enumerate(devices):
        if dev["max_input_channels"] < 1:
            continue  # 入力チャンネルなしはスキップ
        name     = dev["name"]
        is_sys   = any(kw in name.lower() for kw in system_keywords)
        result.append({
            "id":               i,
            "name":             name,
            "is_system_audio":  is_sys,
            "channels":         dev["max_input_channels"],
        })
    return jsonify(result), 200


# ══════════════════════════════════════════════════
#  設定 API
# ══════════════════════════════════════════════════

# PyInstaller frozen 対応
if getattr(sys, "frozen", False):
    ENV_PATH = Path(sys.executable).resolve().parent / ".env"
else:
    ENV_PATH = Path(__file__).resolve().parent / ".env"
# 起動時に .env が存在しない場合は空ファイルを作成
if not ENV_PATH.exists():
    ENV_PATH.write_text("AI_MODE=personal\nRECORDING_SOURCE=mic\n", encoding="utf-8")
    print(f"[main] .env を新規作成しました")

def _read_env() -> dict:
    """`.env` ファイルを読み込んで辞書で返す"""
    env = {}
    if not ENV_PATH.exists():
        return env
    for line in ENV_PATH.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, _, val = line.partition("=")
        env[key.strip()] = val.strip()
    return env

def _write_env(data: dict) -> None:
    """辞書の内容を `.env` ファイルに書き込む（既存の値をマージ）"""
    existing = _read_env()
    existing.update(data)
    lines = [f"{k}={v}" for k, v in existing.items()]
    ENV_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")


@app.route("/api/settings", methods=["GET"])
def settings_get():
    """
    現在の設定を返す（APIキーはマスキング）。
    """
    env = _read_env()

    def _mask(key_name: str) -> dict:
        val    = env.get(key_name, "")
        masked = ("*" * (len(val) - 4) + val[-4:]) if len(val) > 4 else val
        return {"value": masked, "has_key": bool(val)}

    return jsonify({
        "ai_mode":             env.get("AI_MODE", "ollama"),
        "groq_api_key":        _mask("GROQ_API_KEY")["value"],
        "has_groq_key":        _mask("GROQ_API_KEY")["has_key"],
        "openai_api_key":      _mask("OPENAI_API_KEY")["value"],
        "has_openai_key":      _mask("OPENAI_API_KEY")["has_key"],
        "gemini_api_key":      _mask("GEMINI_API_KEY")["value"],
        "has_gemini_key":      _mask("GEMINI_API_KEY")["has_key"],
        "anthropic_api_key":   _mask("ANTHROPIC_API_KEY")["value"],
        "has_anthropic_key":   _mask("ANTHROPIC_API_KEY")["has_key"],
        "recording_source":    env.get("RECORDING_SOURCE", "mic"),
        "recording_device_id": env.get("RECORDING_DEVICE_ID", ""),
        "ollama_model":        env.get("OLLAMA_MODEL", "gemma3:27b"),
    }), 200


@app.route("/api/settings", methods=["POST"])
def settings_save():
    """
    設定を `.env` ファイルに保存する。
    Request JSON: {"ai_mode": str, "groq_api_key": str, ...}
    Response: {"status": "saved"}
    """
    data       = request.get_json(silent=True) or {}
    ai_mode    = data.get("ai_mode", "ollama")
    rec_source = data.get("recording_source", "mic")
    device_id  = data.get("recording_device_id", "")

    save_data = {"AI_MODE": ai_mode, "RECORDING_SOURCE": rec_source}

    # Ollamaモデル
    ollama_model = data.get("ollama_model", "").strip()
    if ollama_model:
        save_data["OLLAMA_MODEL"] = ollama_model

    # 各APIキー（マスク済み・空の場合は上書きしない）
    for key_field, env_key in [
        ("groq_api_key",      "GROQ_API_KEY"),
        ("openai_api_key",    "OPENAI_API_KEY"),
        ("gemini_api_key",    "GEMINI_API_KEY"),
        ("anthropic_api_key", "ANTHROPIC_API_KEY"),
    ]:
        val = data.get(key_field, "").strip()
        if val and not val.startswith("*"):
            save_data[env_key] = val

    # デバイスID
    if device_id and str(device_id).isdigit():
        save_data["RECORDING_DEVICE_ID"] = str(device_id)

    try:
        _write_env(save_data)
        # Ollamaモードに切り替えた場合はモデルをバックグラウンドでプリロード
        if ai_mode == "ollama":
            import platform as _platform
            is_win = _platform.system() == "Windows"
            is_arm = _platform.machine() in ("arm64", "aarch64")
            if is_win or is_arm:
                from app.services.transcriber import preload_model
                preload_model()
        return jsonify({"status": "saved"}), 200
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)}), 500




@app.route("/api/llm/providers", methods=["GET"])
def llm_providers():
    """llm/配下のプロバイダー一覧を返す（設定画面の選択肢生成用）"""
    from app.services.llm import get_all_providers
    return jsonify(get_all_providers()), 200


_ollama_pull_status = {"status": "idle", "model": "", "message": "", "progress": 0}

RECOMMENDED_MODELS = [
    {
        "name": "gemma4:e4b", "size": "約5GB", "desc": "最新世代・高速・推奨",
        "min_ram": "8GB",  "min_vram": "6GB",
        "note": "VRAM 8GB以上のGPUで全量GPU動作。現環境で最速の選択肢。",
        "recommend": "◎",
    },
    {
        "name": "gemma4:26b-a4b", "size": "約17GB", "desc": "最新世代・高性能MoE",
        "min_ram": "32GB", "min_vram": "なし（RAM補完可）",
        "note": "Mixture-of-Experts構造で実効パラメータは4B相当。gemma3:27bと同程度の速度。",
        "recommend": "○",
    },
    {
        "name": "gemma3:27b", "size": "約17GB", "desc": "高品質・日本語得意",
        "min_ram": "32GB", "min_vram": "なし（RAM補完可）",
        "note": "日本語精度が高く安定している。VRAM不足時はCPU補完で低速になる。",
        "recommend": "○",
    },
    {
        "name": "gemma3:12b", "size": "約8GB", "desc": "軽量・バランス型",
        "min_ram": "16GB", "min_vram": "8GB",
        "note": "VRAM 8GBに収まるため高速動作。精度はやや劣るが日常用途に十分。",
        "recommend": "◎",
    },
    {
        "name": "gemma4:31b", "size": "約20GB", "desc": "最新世代・最高性能Dense",
        "min_ram": "32GB", "min_vram": "なし（RAM補完可）",
        "note": "Dense構造で全パラメータが常時稼働。VRAM不足時は大幅に低速。M5 Max以降推奨。",
        "recommend": "△",
    },
    {
        "name": "llama3.3:70b", "size": "約40GB", "desc": "超高性能（高スペック専用）",
        "min_ram": "64GB", "min_vram": "なし（RAM必須）",
        "note": "64GB RAM環境でのみ動作可能。処理は非常に低速。M5 Max 64GB以上を推奨。",
        "recommend": "△",
    },
]

@app.route("/api/ollama/models", methods=["GET"])
def ollama_models():
    """インストール済みモデル一覧と推奨モデルを返す"""
    import urllib.request, json as _json2
    installed = []
    try:
        req = urllib.request.Request("http://127.0.0.1:11434/api/tags")
        with urllib.request.urlopen(req, timeout=3) as res:
            data = _json2.loads(res.read().decode("utf-8"))
            installed = [m["name"] for m in data.get("models", [])]
    except Exception:
        pass
    recommended = []
    for m in RECOMMENDED_MODELS:
        is_installed = any(m["name"] == i or m["name"] == i.split(":")[0] for i in installed)
        recommended.append({**m, "installed": is_installed})
    return jsonify({"status": "ok", "models": installed, "recommended": recommended}), 200


@app.route("/api/ollama/pull", methods=["POST"])
def ollama_pull():
    """指定モデルをバックグラウンドでpullする"""
    global _ollama_pull_status
    data  = request.get_json(silent=True) or {}
    model = data.get("model", "").strip()
    if not model:
        return jsonify({"status": "error", "message": "モデル名が必要です"}), 400
    if _ollama_pull_status.get("status") == "running":
        return jsonify({"status": "error", "message": "既にダウンロード中です"}), 400
    _ollama_pull_status = {"status": "running", "model": model, "message": "ダウンロード開始…", "progress": 0}

    def _pull():
        global _ollama_pull_status
        import subprocess, re as _re2
        try:
            proc = subprocess.Popen(
                ["ollama", "pull", model],
                stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                text=True, bufsize=1
            )
            for line in proc.stdout:
                line = line.strip()
                if not line: continue
                _ollama_pull_status["message"] = line
                m2 = _re2.search(r"(\d+)%", line)
                if m2:
                    _ollama_pull_status["progress"] = int(m2.group(1))
            proc.wait()
            if proc.returncode == 0:
                _ollama_pull_status = {"status": "done", "model": model, "message": f"✅ {model} のインストールが完了しました", "progress": 100}
            else:
                _ollama_pull_status = {"status": "error", "model": model, "message": "ダウンロードに失敗しました", "progress": 0}
        except Exception as e:
            _ollama_pull_status = {"status": "error", "model": model, "message": str(e), "progress": 0}

    import threading as _th2
    _th2.Thread(target=_pull, daemon=True).start()
    return jsonify({"status": "started"}), 200


@app.route("/api/ollama/pull/status", methods=["GET"])
def ollama_pull_status_api():
    """pullの進捗を返す"""
    return jsonify(_ollama_pull_status), 200

@app.route("/api/ollama/status", methods=["GET"])
def ollama_status():
    """Ollamaの起動状態を確認する"""
    import urllib.request
    try:
        urllib.request.urlopen("http://127.0.0.1:11434", timeout=2)
        return jsonify({"status": "running"}), 200
    except Exception:
        return jsonify({"status": "not_running"}), 200



@app.route("/api/clean", methods=["POST"])
def clean_only():
    """文字起こし済みのレコードに対してクリーニングのみ再実行する"""
    global _transcribe_progress, _transcribe_start_time
    import time as _time

    data      = request.get_json(silent=True) or {}
    record_id = data.get("record_id")

    if not record_id:
        return jsonify({"status": "error", "message": "record_id が必要です"}), 400

    record = get_recording(record_id)
    if not record:
        return jsonify({"status": "error", "message": "データが見つかりません"}), 404

    if not record.get("transcript"):
        return jsonify({"status": "error", "message": "文字起こしデータがありません"}), 400

    _transcribe_start_time = _time.time()
    _transcribe_progress   = {"status": "running", "step": "cleaning", "message": "🧹 クリーニング開始…", "progress": 5, "record_id": record_id}

    def on_progress(i, total):
        global _transcribe_progress
        pct = 5 + int(90 * i / max(total, 1))
        _transcribe_progress = {"status": "running", "step": "cleaning", "message": f"🧹 クリーニング中… {i}/{total} ブロック", "progress": pct, "record_id": record_id}

    try:
        env     = _read_env()
        ai_mode = env.get("AI_MODE", "ollama")

        import app.services.transcriber as _tr
        if ai_mode == "ollama":
            cleaned = _tr._clean_transcript_ollama(record["transcript"], progress_cb=on_progress)
        else:
            # クラウドAPIモード：Groq Whisperは不要・LLMのみ使用
            api_key = env.get("GROQ_API_KEY", "").strip()
            if not api_key:
                return jsonify({"status": "error", "message": "Groq APIキーが設定されていません"}), 400
            from groq import Groq
            client  = Groq(api_key=api_key)
            cleaned = _tr._clean_transcript_groq(client, record["transcript"])

        import app.database as _db
        _db.update_cleaned_transcript(record_id, cleaned)
        _transcribe_progress = {"status": "done", "step": "done", "message": "✅ クリーニング完了！", "progress": 100, "record_id": record_id}
        return jsonify({"status": "done", "cleaned_transcript": cleaned}), 200

    except Exception as e:
        from app.services.transcriber import _parse_api_error
        err_msg = _parse_api_error(e, "クリーニング")
        update_cleaning_error(record_id, err_msg)
        _transcribe_progress = {"status": "error", "step": "error", "message": f"⚠️ {err_msg}", "progress": 0, "record_id": record_id}
        return jsonify({"status": "error", "message": err_msg}), 500


@app.route("/api/summarize", methods=["POST"])
def summarize_only():
    """文字起こし済みのレコードに対して要約のみ再実行する"""
    global _transcribe_progress, _transcribe_start_time
    import time as _time

    data      = request.get_json(silent=True) or {}
    record_id = data.get("record_id")

    if not record_id:
        return jsonify({"status": "error", "message": "record_id が必要です"}), 400

    record = get_recording(record_id)
    if not record:
        return jsonify({"status": "error", "message": "データが見つかりません"}), 404

    if not record.get("transcript"):
        return jsonify({"status": "error", "message": "文字起こしデータがありません"}), 400

    _transcribe_start_time = _time.time()
    _transcribe_progress   = {"status": "running", "step": "summarizing", "message": "✨ 要約開始…", "progress": 5, "record_id": record_id}

    def on_progress(i, total, final=False):
        global _transcribe_progress
        if final:
            _transcribe_progress = {"status": "running", "step": "summarizing", "message": "✨ 統合要約中…", "progress": 90, "record_id": record_id}
        else:
            pct = 5 + int(80 * i / max(total, 1))
            _transcribe_progress = {"status": "running", "step": "summarizing", "message": f"✨ 要約中… {i}/{total} ブロック", "progress": pct, "record_id": record_id}

    try:
        extra_prompt = data.get("extra_prompt", "").strip()
        env     = _read_env()
        ai_mode = env.get("AI_MODE", "ollama")

        # cleaned_transcriptを優先、なければtranscriptにフォールバック
        source_text = record.get("cleaned_transcript") or record.get("transcript", "")

        if ai_mode == "ollama":
            from app.services.transcriber import _summarize_ollama
            ai_summary = _summarize_ollama(source_text, extra_prompt, progress_cb=on_progress)
        else:
            from app.services.transcriber import _summarize_cloud
            ai_summary = _summarize_cloud(source_text, extra_prompt, ai_mode, env)

        # エラーメッセージが返された場合はエラーとして扱う
        if ai_summary.startswith("（") and ai_summary.endswith("）"):
            err_msg = ai_summary[1:-1]
            # summary_errorに保存（ai_summaryは更新しない）
            update_summary_error(record_id, err_msg)
            _transcribe_progress = {"status": "error", "step": "error", "message": f"⚠️ {err_msg}", "progress": 0, "record_id": record_id}
            return jsonify({"status": "error", "message": err_msg}), 500

        # 成功時：ai_summaryを保存しsummary_errorをクリア
        update_transcript_and_summary(record_id, record["transcript"], ai_summary)
        _transcribe_progress = {"status": "done", "step": "done", "message": "✅ 要約完了！", "progress": 100, "record_id": record_id}
        return jsonify({"status": "done", "ai_summary": ai_summary}), 200
    except Exception as e:
        from app.services.transcriber import _parse_api_error
        env2     = _read_env()
        ai_mode2 = env2.get("AI_MODE", "ollama")
        from app.services.llm import get_provider as _gp
        _prov  = _gp(ai_mode2)
        _pname = getattr(_prov, "DISPLAY_NAME", ai_mode2) if _prov else ai_mode2
        err_msg = _parse_api_error(e, _pname)
        update_summary_error(record_id, err_msg)
        _transcribe_progress = {"status": "error", "step": "error", "message": f"⚠️ {err_msg}", "progress": 0, "record_id": record_id}
        return jsonify({"status": "error", "message": err_msg}), 500

@app.route("/api/model/status", methods=["GET"])
def model_status():
    """
    モデルのダウンロード・ロード状態を返す。
    Response: {"status": "idle"|"downloading"|"ready"|"error", "error": str}
    """
    try:
        from app.services.transcriber import get_model_status
        return jsonify(get_model_status()), 200
    except Exception as e:
        return jsonify({"status": "idle", "error": str(e)}), 200


# ══════════════════════════════════════════════════
#  相談室 API
# ══════════════════════════════════════════════════

@app.route("/api/personas", methods=["GET"])
def personas_get():
    """ペルソナ設定一覧を返す"""
    return jsonify(get_all_persona_settings()), 200


@app.route("/api/personas", methods=["POST"])
def personas_create():
    """新しいペルソナを作成する"""
    data = request.get_json(silent=True) or {}
    name   = data.get("persona_name", "").strip()
    role   = data.get("role", "").strip()
    prompt = data.get("system_prompt", "").strip()
    if not name:
        return jsonify({"status": "error", "message": "名前は必須です"}), 400
    result = create_persona(name, role, prompt)
    if result.get("status") == "ok":
        # 新規ペルソナを「グループなし」に追加
        add_persona_to_default_group(name)
    return jsonify(result), 200


@app.route("/api/personas/<string:persona_name>", methods=["PUT"])
def personas_update(persona_name: str):
    """ペルソナを更新する"""
    data     = request.get_json(silent=True) or {}
    new_name = data.get("persona_name", persona_name).strip()
    role     = data.get("role", "").strip()
    prompt   = data.get("system_prompt", "").strip()
    enabled  = bool(data.get("enabled", True))
    # 後方互換：enabled のみの更新も受け付ける
    if "enabled" in data and len(data) == 1:
        return jsonify(update_persona_enabled(persona_name, enabled)), 200
    result = update_persona(persona_name, new_name, role, prompt, enabled)
    # ペルソナ名が変わった場合はグループのメンバー名も自動更新
    if result.get("status") == "ok" and new_name != persona_name:
        with get_connection() as conn:
            conn.execute(
                "UPDATE persona_group_members SET persona_name=? WHERE persona_name=?",
                (new_name, persona_name)
            )
            conn.commit()
    return jsonify(result), 200


@app.route("/api/personas/<string:persona_name>", methods=["DELETE"])
def personas_delete(persona_name: str):
    """ペルソナを削除する"""
    remove_persona_from_all_groups(persona_name)
    return jsonify(delete_persona(persona_name)), 200


@app.route("/api/persona_groups", methods=["GET"])
def persona_groups_get():
    """ペルソナグループ一覧をメンバー付きで返す"""
    return jsonify(get_all_persona_groups()), 200


@app.route("/api/persona_groups", methods=["POST"])
def persona_groups_create():
    """グループを新規作成する"""
    data = request.get_json(silent=True) or {}
    name = data.get("name", "").strip()
    if not name:
        return jsonify({"status": "error", "message": "グループ名は必須です"}), 400
    return jsonify(create_persona_group(name)), 200


@app.route("/api/persona_groups/<int:group_id>", methods=["PUT"])
def persona_groups_update(group_id: int):
    """グループ名とメンバーを更新する"""
    data    = request.get_json(silent=True) or {}
    name    = data.get("name", "").strip()
    members = data.get("members", [])
    if not name:
        return jsonify({"status": "error", "message": "グループ名は必須です"}), 400
    return jsonify(update_persona_group(group_id, name, members)), 200


@app.route("/api/persona_groups/<int:group_id>", methods=["DELETE"])
def persona_groups_delete(group_id: int):
    """グループを削除する（メンバーは「グループなし」に移動）"""
    return jsonify(delete_persona_group(group_id)), 200


@app.route("/api/council/ask", methods=["POST"])
def council_ask():
    """
    相談を受け付けてSSEでペルソナ回答をストリーミングする。
    Request JSON: {"question": str, "category_id": int, "personas": [str, ...]}
    SSE events:
      {"type": "start",    "persona": str}
      {"type": "chunk",    "persona": str, "chunk": str}
      {"type": "done",     "persona": str, "answer": str}
      {"type": "finished"}
      {"type": "error",    "message": str}
    """
    import json as _json
    import urllib.request as _ureq
    from app.services.transcriber import _get_ollama_model
    from flask import Response, stream_with_context

    data          = request.get_json(silent=True) or {}
    question      = data.get("question", "").strip()
    # 複数カテゴリ対応
    cat_ids_raw   = data.get("category_ids")
    if cat_ids_raw:
        category_ids  = [int(c) for c in cat_ids_raw if c]
        category_id   = category_ids[0]
    else:
        category_id   = int(data.get("category_id", 1) or 1)
        category_ids  = [category_id]
    persona_names = data.get("personas", [])
    context_keys  = set(data.get("context_keys", None) or [])
    council_history = data.get("council_history", [])  # 連続相談用会話履歴
    use_all_context = data.get("context_keys") is None

    if not question:
        return jsonify({"status": "error", "message": "質問を入力してください"}), 400
    if not persona_names:
        return jsonify({"status": "error", "message": "参加ペルソナがいません"}), 400

    # ── コンテキスト収集（複数カテゴリ対応）──
    ctx = get_context_for_categories(category_ids)

    def build_context_text():
        if not use_all_context and not context_keys:
            return ""

        # カテゴリIDと名前のマップを作成
        all_cats = get_all_categories()
        cat_name_map = {c['id']: c['name'] for c in all_cats}
        multi_cat = len(category_ids) > 1  # 複数カテゴリ選択中かどうか

        def cat_label(item):
            if not multi_cat:
                return ""
            cid = item.get('category_id')
            return f"({cat_name_map.get(cid, f'カテゴリ{cid}')})" if cid else ""

        parts = []
        for r in ctx.get("recordings", []):
            key = f"rec_{r['id']}"
            if not use_all_context and key not in context_keys:
                continue
            # ai_summaryを優先、なければcleaned_transcript（5000文字上限）にフォールバック
            ref_text = r.get("ai_summary") or r.get("cleaned_transcript") or ""
            if ref_text:
                parts.append(f"[録音:{r['title']}{cat_label(r)}]\n{ref_text[:5000]}")
        for m in ctx.get("memos", []):
            key = f"memo_{m['id']}"
            if not use_all_context and key not in context_keys:
                continue
            parts.append(f"[メモ:{m['title']}{cat_label(m)}]\n{m['body'][:5000]}")
        for f in ctx.get("files", []):
            key = f"file_{f['id']}"
            if not use_all_context and key not in context_keys:
                continue
            if f.get("summary"):
                parts.append(f"[ファイル:{f['original_name']}{cat_label(f)}]\n{f['summary'][:5000]}")
        for s in ctx.get("sessions", []):
            if s.get("final_decision"):
                label = cat_label(s)
                parts.append(f"[過去の決定{label}] 質問:{s['question']} → {s['final_decision']}")
        return "\n\n".join(parts)

    context_text = build_context_text()

    # ── 関連性フィルタリング：質問と無関係なコンテキストを除外 ──
    def filter_relevant_context(question: str, raw_context: str, env: dict) -> str:
        """
        Ollamaに各コンテキスト項目と質問の関連性を判定させ、
        関連ありと判定されたもののみ返す。
        コンテキストが空、またはOllamaが使えない場合はそのまま返す。
        """
        if not raw_context.strip():
            return ""

        # 項目単位に分割（[タグ:タイトル] で始まるブロック）
        import re as _re
        blocks = _re.split(r'\n\n(?=\[)', raw_context.strip())
        if len(blocks) <= 1:
            # 分割できない場合はそのまま
            return raw_context

        ai_mode = env.get("AI_MODE", "ollama")
        model   = _get_ollama_model()

        relevant_blocks = []
        for block in blocks:
            first_line = block.split('\n')[0]
            snippet    = block[:300]

            judge_prompt = (
                f"以下の「参考情報」は、「質問」に回答する際に役立つ情報ですか？\n\n"
                f"質問：{question}\n\n"
                f"参考情報：{snippet}\n\n"
                f"「はい」か「いいえ」の一単語のみで答えてください。"
            )

            try:
                from app.services.llm import get_provider, is_local_mode
                provider = get_provider(ai_mode)

                if is_local_mode(ai_mode):
                    import json as _j, urllib.request as _ur
                    payload = _j.dumps({
                        "model": model,
                        "messages": [{"role": "user", "content": judge_prompt}],
                        "stream": False,
                        "options": {"temperature": 0, "num_predict": 10, "num_ctx": 512},
                    }).encode("utf-8")
                    req = _ur.Request(
                        "http://127.0.0.1:11434/api/chat",
                        data=payload,
                        headers={"Content-Type": "application/json"},
                        method="POST",
                    )
                    with _ur.urlopen(req, timeout=15) as res:
                        result = _j.loads(res.read())
                    answer = result.get("message", {}).get("content", "").strip().lower()
                else:
                    api_key_name = getattr(provider, "API_KEY_NAME", "")
                    api_key      = env.get(api_key_name, "").strip() if api_key_name else ""
                    if not api_key:
                        relevant_blocks.append(block)
                        continue
                    messages = [{"role": "user", "content": judge_prompt}]
                    answer   = provider.chat(messages, api_key=api_key, options={"temperature": 0, "max_tokens": 5}).lower()

                if "はい" in answer or "yes" in answer:
                    relevant_blocks.append(block)
                elif not answer:
                    relevant_blocks.append(block)

            except Exception:
                relevant_blocks.append(block)

        return "\n\n".join(relevant_blocks)

    # ── ペルソナ定義：DBから動的読み込み ──
    all_personas = get_all_persona_settings()
    persona_db_map = {p['persona_name']: p for p in all_personas}

    # 共通フォーマット指示は廃止 → 各ペルソナのプロンプトに個別定義
    PERSONA_FORMAT_SUFFIX = ""  # 互換性のため残すが空にする

    # 全ペルソナ共通の情報読み取り指示（先頭に付与・簡潔版）
    PERSONA_COMMON_PREFIX = """【役割と指示】
あなたは以下に定義された専門家として回答する。

【入力データの扱い】
- 「相談者の基本情報」→ 回答の前提として必ず反映する
- 「参考情報」→ 相談内容と関連する場合のみ使う。無関係なら無視
- 「相談内容」→ これに答える

"""

    # ペルソナ別 temperature（創造性）設定
    # 高いほど個性的・低いほど論理的・一貫
    PERSONA_TEMPERATURE = {
        '戦略家':       0.6,  # 断定的・一貫性重視
        'リスク管理者': 0.4,  # 慎重・論理的
        'アナリスト':   0.3,  # データ重視・客観的
        'クリエイター': 0.95, # 斬新・自由な発想
        '法務・コンプラ': 0.3, # 厳格・一貫性重視
        'ユーザー視点': 0.7,  # 感情・直感を重視
        'クレーマー':   0.8,  # 攻撃的・予測不能
        'マーケッター': 0.75, # 市場感覚・柔軟
    }
    DEFAULT_TEMPERATURE = 0.7

    def generate():
        ollama_res = None  # クライアント切断時に閉じるための参照
        try:
            env     = _read_env()
            ai_mode = env.get("AI_MODE", "ollama")

            from app.services.llm import get_provider, is_local_mode
            provider = get_provider(ai_mode)
            if provider is None:
                yield f"data: {_json.dumps({'type': 'error', 'message': f'不明なAIモード: {ai_mode}'}, ensure_ascii=False)}\n\n"
                return

            # クラウドモード時のAPIキー取得
            api_key = ""
            if not is_local_mode(ai_mode):
                api_key_name = getattr(provider, "API_KEY_NAME", "")
                api_key      = env.get(api_key_name, "").strip() if api_key_name else ""
                if not api_key:
                    yield f"data: {_json.dumps({'type': 'error', 'message': f'{provider.DISPLAY_NAME} のAPIキーが設定されていません'}, ensure_ascii=False)}\n\n"
                    return

            # ── 関連性フィルタリングを事前に実行（全ペルソナ共通） ──
            filtered_context = filter_relevant_context(question, context_text, env)

            for persona_name in persona_names:
                p_db = persona_db_map.get(persona_name)
                if p_db and p_db.get('system_prompt'):
                    base_prompt = p_db['system_prompt'] + PERSONA_FORMAT_SUFFIX
                else:
                    base_prompt = f"あなたは{persona_name}の専門家です。自分の立場から400字以内・日本語で答えてください。"
                system_content = PERSONA_COMMON_PREFIX + base_prompt

                parts_msg = []
                if filtered_context:
                    parts_msg.append(
                        "■ 参考情報（相談内容に関連すると判断されたデータ）\n"
                        f"{filtered_context}"
                    )

                try:
                    high_rated = get_highly_rated_answers(persona_name, limit=2)
                    if high_rated:
                        ref_parts = [f"■ {persona_name}の過去の高評価回答（参考）\n以下はこのペルソナが過去に高く評価された回答例です。回答スタイルの参考にしてください。"]
                        for hr in high_rated:
                            ref_parts.append(f"【質問】{hr['question']}\n【回答(★{hr['rating']})】{hr['answer'][:300]}")
                        parts_msg.append("\n\n".join(ref_parts))
                except Exception:
                    pass

                parts_msg.append(f"■ 相談内容\n{question}")
                user_content = "\n\n" + "\n\n".join(parts_msg)

                yield f"data: {_json.dumps({'type': 'start', 'persona': persona_name}, ensure_ascii=False)}\n\n"

                full_answer = ""
                try:
                    messages = [{"role": "system", "content": system_content}]
                    for h in council_history:
                        role     = h.get("role", "user")
                        hcontent = h.get("content", "")
                        if role in ("user", "assistant") and hcontent:
                            messages.append({"role": role, "content": hcontent})
                    messages.append({"role": "user", "content": user_content})

                    from app.services.transcriber import _get_summary_num_ctx as _gsnc
                    options = {
                        "temperature":    PERSONA_TEMPERATURE.get(persona_name, DEFAULT_TEMPERATURE),
                        "num_predict":    2048,
                        "num_ctx":        _gsnc(),
                        "repeat_penalty": 1.2,
                    }

                    if is_local_mode(ai_mode):
                        for chunk in provider.chat_stream(messages, options=options):
                            full_answer += chunk
                            yield f"data: {_json.dumps({'type': 'chunk', 'persona': persona_name, 'chunk': chunk}, ensure_ascii=False)}\n\n"
                    else:
                        for chunk in provider.chat_stream(messages, api_key=api_key, options=options):
                            full_answer += chunk
                            yield f"data: {_json.dumps({'type': 'chunk', 'persona': persona_name, 'chunk': chunk}, ensure_ascii=False)}\n\n"

                except Exception as e:
                    err_msg = str(e)
                    print(f"[council] ペルソナ '{persona_name}' エラー: {err_msg}")
                    yield f"data: {_json.dumps({'type': 'error', 'message': err_msg}, ensure_ascii=False)}\n\n"
                    full_answer = ""
                    continue

                yield f"data: {_json.dumps({'type': 'done', 'persona': persona_name, 'answer': full_answer}, ensure_ascii=False)}\n\n"

            yield f"data: {_json.dumps({'type': 'finished'}, ensure_ascii=False)}\n\n"

        except GeneratorExit:
            print("[council] クライアント切断を検知。処理を中断します。")
            try:
                if ollama_res:
                    ollama_res.close()
            except Exception:
                pass
            return

        except Exception as e:
            yield f"data: {_json.dumps({'type': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.route("/api/council/save", methods=["POST"])
def council_save():
    data           = request.get_json(silent=True) or {}
    question       = data.get("question", "").strip()
    category_id    = int(data.get("category_id", 1) or 1)
    final_decision = data.get("final_decision", "").strip()
    adopted_list   = data.get("adopted", [])

    if not question:
        return jsonify({"status": "error", "message": "質問は必須です"}), 400

    session_id = create_council_session(question, category_id, final_decision)
    print(f"[council] セッション保存: id={session_id} adopted={len(adopted_list)}件")
    for a in adopted_list:
        pname  = a.get("persona_name", "").strip()
        answer = a.get("answer", "").strip()
        print(f"[council] 採用回答: persona={pname} answer_len={len(answer)}")
        if pname and answer:
            create_council_adopted(session_id, pname, answer)

    return jsonify({"status": "ok", "session_id": session_id}), 200


@app.route("/api/recordings/<int:record_id>/trim", methods=["POST"])
def recording_trim(record_id: int):
    """指定した開始・終了時間で録音をトリミングする"""
    record = get_recording(record_id)
    if not record:
        return jsonify({"status": "error", "message": "データが見つかりません"}), 404

    wav_path = UPLOADS_DIR / record["wav_file"]
    if not wav_path.exists():
        return jsonify({"status": "error", "message": "音声ファイルが見つかりません"}), 404

    data     = request.get_json(silent=True) or {}
    start    = float(data.get("start", 0))
    end      = float(data.get("end", 0))
    duration = end - start

    if duration <= 0:
        return jsonify({"status": "error", "message": "開始・終了位置が不正です"}), 400

    try:
        import subprocess, shutil, uuid
        tmp_path = UPLOADS_DIR / f"trim_tmp_{uuid.uuid4().hex}.wav"
        print(f"[trim] 実行: ffmpeg -ss {start} -t {duration} {wav_path}")
        result   = subprocess.run([
            "ffmpeg", "-y",
            "-i", str(wav_path),
            "-ss", str(start),
            "-t",  str(duration),
            "-c",  "copy",
            str(tmp_path)
        ], capture_output=True, timeout=120)

        if result.returncode != 0:
            err_msg = result.stderr.decode('utf-8', errors='replace')[-500:]
            print(f"[trim] ffmpeg error: {err_msg}")
            return jsonify({"status": "error", "message": f"ffmpegエラー: {err_msg}"}), 500

        shutil.move(str(tmp_path), str(wav_path))
        print(f"[trim] 完了: {wav_path}")
        return jsonify({"status": "ok"}), 200

    except FileNotFoundError as e:
        print(f"[trim] FileNotFoundError: {e}")
        return jsonify({"status": "error", "message": "ffmpegがインストールされていません"}), 500
    except Exception as e:
        import traceback
        print(f"[trim] 例外: {traceback.format_exc()}")
        return jsonify({"status": "error", "message": str(e)}), 500


@app.route("/api/council/sessions/<int:session_id>", methods=["DELETE"])
def council_session_delete(session_id: int):
    """相談セッションを削除する"""
    return jsonify(delete_council_session(session_id)), 200


@app.route("/api/council/sessions/<int:session_id>", methods=["PUT"])
def council_session_update(session_id: int):
    """既存セッションの最終判断メモと採用回答を更新する"""
    data           = request.get_json(silent=True) or {}
    final_decision = data.get("final_decision", "").strip()
    adopted_list   = data.get("adopted", [])

    update_council_session_decision(session_id, final_decision)

    # 採用回答がある場合は既存を削除して再登録（上書き）
    if adopted_list:
        with get_connection() as conn:
            conn.execute(
                "DELETE FROM council_adopted WHERE session_id=?", (session_id,)
            )
            conn.commit()
        for a in adopted_list:
            pname  = a.get("persona_name", "").strip()
            answer = a.get("answer", "").strip()
            if pname and answer:
                create_council_adopted(session_id, pname, answer)

    return jsonify({"status": "ok", "session_id": session_id}), 200


@app.route("/api/personas/<string:persona_name>/icon", methods=["POST"])
def personas_icon_upload(persona_name: str):
    """ペルソナのアイコン画像をアップロードする（variant: wait or think）"""
    if "file" not in request.files:
        return jsonify({"status": "error", "message": "ファイルが選択されていません"}), 400

    file    = request.files["file"]
    variant = request.form.get("variant", "wait")  # 'wait' or 'think'
    if variant not in ("wait", "think"):
        variant = "wait"

    ext = Path(file.filename).suffix.lower()
    if ext not in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        return jsonify({"status": "error", "message": "PNG・JPG・WebP・GIF のみ対応しています"}), 400

    file.seek(0, 2); size = file.tell(); file.seek(0)
    if size > 2 * 1024 * 1024:
        return jsonify({"status": "error", "message": "ファイルサイズは2MB以内にしてください"}), 400

    icon_dir = RESOURCE_DIR / "app" / "static" / "img" / "personas" / "custom"
    icon_dir.mkdir(parents=True, exist_ok=True)

    import hashlib
    safe_name = hashlib.md5(persona_name.encode()).hexdigest()
    save_path = icon_dir / f"{safe_name}_{variant}.png"

    try:
        from PIL import Image as _PILImage
        img = _PILImage.open(file.stream).convert("RGBA")
        img = img.resize((256, 256), _PILImage.LANCZOS)
        img.save(str(save_path), "PNG")
    except Exception:
        file.seek(0)
        file.save(str(save_path))

    url = f"/static/img/personas/custom/{safe_name}_{variant}.png"
    return jsonify({"status": "ok", "variant": variant, "path": url}), 200


@app.route("/api/personas/<string:persona_name>/icon", methods=["DELETE"])
def personas_icon_delete(persona_name: str):
    """ペルソナのカスタムアイコンを削除する（variant指定可）"""
    import hashlib
    safe_name = hashlib.md5(persona_name.encode()).hexdigest()
    variant   = request.args.get("variant", "all")
    icon_dir  = RESOURCE_DIR / "app" / "static" / "img" / "personas" / "custom"

    if variant == "all":
        for v in ("wait", "think"):
            p = icon_dir / f"{safe_name}_{v}.png"
            if p.exists(): p.unlink()
    else:
        p = icon_dir / f"{safe_name}_{variant}.png"
        if p.exists(): p.unlink()

    return jsonify({"status": "ok"}), 200


@app.route("/api/personas/<string:persona_name>/icon", methods=["GET"])
def personas_icon_get(persona_name: str):
    """ペルソナのカスタムアイコンURL（wait/think両方）を返す"""
    import hashlib
    safe_name = hashlib.md5(persona_name.encode()).hexdigest()
    icon_dir  = RESOURCE_DIR / "app" / "static" / "img" / "personas" / "custom"
    result    = {}
    for v in ("wait", "think"):
        p = icon_dir / f"{safe_name}_{v}.png"
        result[v] = f"/static/img/personas/custom/{safe_name}_{v}.png" if p.exists() else None
    return jsonify(result), 200


@app.route("/api/council/sessions", methods=["GET"])
def council_sessions_get():
    """相談履歴を検索・フィルターして返す"""
    category_id = request.args.get("category_id", type=int)
    keyword     = request.args.get("keyword", "").strip()
    date_from   = request.args.get("date_from", "").strip()
    date_to     = request.args.get("date_to", "").strip()
    sessions    = get_council_sessions(
        category_id=category_id,
        keyword=keyword,
        date_from=date_from,
        date_to=date_to,
    )
    return jsonify(sessions), 200


@app.route("/api/council/adopted/<int:adopted_id>/rating", methods=["PUT"])
def council_adopted_rating(adopted_id: int):
    """採用回答の評価を更新する"""
    data   = request.get_json(silent=True) or {}
    rating = int(data.get("rating", 0))
    return jsonify(update_adopted_rating(adopted_id, rating)), 200


# ══════════════════════════════════════════════════
#  起動
# ══════════════════════════════════════════════════


def _start_ollama():
    """COVEと一緒にOllamaを起動する（ビジネス用モード時）"""
    import subprocess
    import urllib.request
    # すでに起動しているか確認
    try:
        urllib.request.urlopen("http://127.0.0.1:11434", timeout=2)
        print("[main] Ollamaはすでに起動しています")
        return
    except Exception:
        pass
    # Ollamaを起動（最後のリクエストから5分後にモデルを自動アンロード）
    try:
        env = os.environ.copy()
        env["OLLAMA_KEEP_ALIVE"] = "5m"  # 5分間使用がなければRAMを解放
        subprocess.Popen(
            ["ollama", "serve"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            env=env,
        )
        print("[main] Ollamaを起動しました（KEEP_ALIVE=5m）")
    except Exception as e:
        print(f"[main] Ollama起動エラー: {e}")

def _open_browser():
    """Flaskの起動を待ってからブラウザを開く"""
    import time
    time.sleep(1.5)
    webbrowser.open("http://127.0.0.1:5001")


def _auto_preload_model():
    """起動時にOllamaモードなら自動でモデルをプリロードする"""
    # Flaskデバッグモードはリローダーが親・子の2プロセスを起動する
    # WERKZEUG_RUN_MAIN=true の子プロセスのみで実行し、2重ロードを防ぐ
    if not getattr(sys, "frozen", False):
        if os.environ.get("WERKZEUG_RUN_MAIN") != "true":
            return
    try:
        env = _read_env()
        if env.get("AI_MODE") == "ollama":
            import platform as _platform
            system = _platform.system()
            if system in ("Windows", "Darwin", "Linux"):
                from app.services.transcriber import preload_model
                preload_model()
                print("[main] 起動時モデルプリロード開始")
    except Exception as e:
        print(f"[main] 起動時プリロードエラー: {e}")

if __name__ == "__main__":
    # PyInstallerで固めた場合はデバッグ無効・自動リロード無効
    is_frozen = getattr(sys, "frozen", False)

    if not is_frozen:
        # 通常の開発時はデバッグモード
        threading.Thread(target=_start_ollama, daemon=True).start()
        threading.Thread(target=_auto_preload_model, daemon=True).start()
        app.run(host="127.0.0.1", port=5001, debug=True)
    else:
        # 配布用：ブラウザを別スレッドで起動してからFlaskを起動
        print("COVE を起動しています...")
        print("ブラウザが開かない場合は http://127.0.0.1:5001 にアクセスしてください")
        threading.Thread(target=_start_ollama, daemon=True).start()
        threading.Thread(target=_auto_preload_model, daemon=True).start()
        threading.Thread(target=_open_browser, daemon=True).start()
        app.run(host="127.0.0.1", port=5001, debug=False, use_reloader=False)
